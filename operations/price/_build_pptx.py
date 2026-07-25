# -*- coding: utf-8 -*-
"""Build price.pptx — paper-curation 운영 비용 분석 deck.
Style matched to the 홍릉포럼 template (16:9, KoPub돋움체, gray body, red/blue accents)."""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TPL = "/Users/jehyunlee/Documents/내노트북/00_Personal/외부활동/260612_홍릉포럼_김은지/260612_이제현_AI for Science 현황.pptx"
HERE = Path(__file__).resolve().parent
FIG = str(HERE / "figures") + os.sep
OUT = str(HERE / "price.pptx")

# palette (template theme)
RED   = RGBColor.from_string("C00000")
BLUE  = RGBColor.from_string("4472C4")
AMBER = RGBColor.from_string("BF8F00")
GRAY  = RGBColor.from_string("808080")
DARK  = RGBColor.from_string("262626")
NAVY  = RGBColor.from_string("1F3864")
GREEN = RGBColor.from_string("548235")
WHITE = RGBColor.from_string("FFFFFF")
LGRAY = RGBColor.from_string("F2F2F2")
BODY  = "KoPub돋움체 Medium"
BOLD  = "KoPub돋움체 Bold"

prs = Presentation(TPL)
EMU = 914400
SW, SH = prs.slide_width, prs.slide_height

# ---- wipe existing slides ----
lst = prs.slides._sldIdLst
for sid in list(lst):
    rId = sid.get(qn('r:id'))
    try: prs.part.drop_rel(rId)
    except Exception: pass
    lst.remove(sid)

BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def _set(run, size, color, bold=False, font=None):
    run.font.size = Pt(size); run.font.color.rgb = color; run.font.bold = bold
    run.font.name = font or (BOLD if bold else BODY)

def tb(s, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2); tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    return box, tf

def para(tf, text, size, color, bold=False, first=False, align=PP_ALIGN.LEFT,
         space_after=4, bullet=None, font=None, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after); p.level = level
    if bullet:
        r = p.add_run(); _set(r, size, bullet[1] if len(bullet)>1 else color, True); r.text = bullet[0]+" "
    r = p.add_run(); _set(r, size, color, bold, font); r.text = text
    return p

def runs(tf, parts, size, first=False, align=PP_ALIGN.LEFT, space_after=4):
    """parts = list of (text, color, bold)."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space_after)
    for (text, color, bold) in parts:
        r = p.add_run(); _set(r, size, color, bold); r.text = text
    return p

def title(s, text, sub=None):
    box, tf = tb(s, 0.5, 0.28, 12.3, 0.85)
    para(tf, text, 25, DARK, True, first=True)
    # red accent bar
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.52), Inches(1.06), Inches(2.1), Inches(0.055))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()
    if sub:
        sbox, stf = tb(s, 0.54, 1.13, 12.2, 0.4)
        para(stf, sub, 12.5, GRAY, False, first=True)

def source(s, text):
    box, tf = tb(s, 0.5, 7.06, 12.3, 0.36)
    para(tf, text, 9, GRAY, False, first=True)

def pic(s, path, l, t, w=None, h=None):
    if not os.path.exists(path): return None
    kw = {}
    if w: kw["width"] = Inches(w)
    if h: kw["height"] = Inches(h)
    return s.shapes.add_picture(path, Inches(l), Inches(t), **kw)

def chip(s, l, t, w, h, fill, text, tcolor=WHITE, size=12, bold=True, line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, text, size, tcolor, bold, first=True, align=PP_ALIGN.CENTER)
    return sh

def table(s, data, l, t, w, h, col_w=None, header_fill=NAVY, sizes=None, aligns=None):
    rows, cols = len(data), len(data[0])
    gt = s.shapes.add_table(rows, cols, Inches(l), Inches(t), Inches(w), Inches(h)).table
    gt.first_row = True; gt.horz_banding = True
    if col_w:
        for i, cw in enumerate(col_w): gt.columns[i].width = Inches(cw)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            c = gt.cell(ri, ci)
            c.margin_left = Pt(5); c.margin_right = Pt(5); c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = c.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            al = (aligns[ci] if aligns else (PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER))
            p.alignment = al
            r = p.add_run(); r.text = str(val)
            sz = sizes[ri] if sizes else (11 if ri == 0 else 10.5)
            if ri == 0:
                _set(r, sz, WHITE, True); c.fill.solid(); c.fill.fore_color.rgb = header_fill
            else:
                _set(r, sz, DARK, False); c.fill.solid()
                c.fill.fore_color.rgb = WHITE if ri % 2 else LGRAY
    return gt

# ============================ SLIDE 1 — title ============================
s = slide()
# accent band
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.55), Inches(13.333), Inches(0.07))
band.fill.solid(); band.fill.fore_color.rgb = RED; band.line.fill.background(); band.shadow.inherit=False
box, tf = tb(s, 0.7, 1.55, 12, 1.0)
para(tf, "paper-curation 운영 비용 분석", 40, DARK, True, first=True)
box2, tf2 = tb(s, 0.72, 2.75, 12, 1.4)
para(tf2, "동작·옵션별 API 비용 실측 + LLM wiki 운영 Trade-off", 18, NAVY, True, first=True, space_after=6)
para(tf2, "예제 논문: arXiv 2412.11427 “Towards Scientific Discovery with Generative AI” · 실측 기반", 13, GRAY, False)
box3, tf3 = tb(s, 0.72, 6.05, 9, 0.5)
para(tf3, "2026. 06. 19.  ·  한국과학기술연구원  ·  이제현", 14, DARK, False, first=True)
pic(s, FIG+"logo_1.png", 0.14, 6.56, w=2.48, h=0.82)
pic(s, FIG+"logo_0.png", 12.26, 6.43, w=1.07, h=1.07)

# ============================ SLIDE 2 — 개요 ============================
s = slide()
title(s, "무엇을 계산했나 — 동작 × 옵션 × 실측", "Zotero Paper Curio → 파이프라인 → Deep/Deeper Research → wiki 운영까지")
box, tf = tb(s, 0.6, 1.45, 7.4, 5.2)
para(tf, "측정 대상", 14, RED, True, first=True, space_after=6)
for t1 in ["Paper Curio (리뷰·일괄·웹배포)", "파이프라인 (분류·연결·인사이트·타임라인)",
           "Deep Research (길이 × 모델 옵션)", "Deeper Research (멀티에이전트·Opus)",
           "Audio Overview (TTS)", "LLM wiki 운영 (큐레이션 vs PDF 쌓기)"]:
    para(tf, t1, 13, DARK, False, bullet=("•", RED), space_after=5)
para(tf, "측정 방법", 14, RED, True, space_after=6)
para(tf, "입력 토큰 = Anthropic count_tokens (과금 0) · 출력 = 실제 생성 usage", 12.5, DARK, False, bullet=("✓", GREEN))
para(tf, "단가 = 2026-06 공식 (Anthropic·OpenAI·Google)", 12.5, DARK, False, bullet=("✓", GREEN))
# 단가 mini table
table(s, [["모델", "In", "Out"],
          ["Haiku 4.5", "$1", "$5"],
          ["Sonnet 4.6", "$3", "$15"],
          ["Opus 4.8", "$5", "$25"],
          ["Gemini 3.5-flash", "$1.5", "$9"],
          ["Gemini embed", "$0.15", "—"]],
      8.25, 1.6, 4.5, 3.5, col_w=[2.5, 1.0, 1.0], sizes=[12,12,12,12,12,12])
box2, tf2 = tb(s, 8.25, 5.25, 4.5, 1.4)
para(tf2, "per 1M tokens · 답변 BYOK(독자) / 큐레이션 운영자 부담", 11, GRAY, False, first=True)
source(s, "단가: Anthropic·OpenAI·Google 공식 페이지(2026-06) · 토큰: count_tokens + 실측")

# ============================ SLIDE 3 — 논문 1편 ============================
s = slide()
title(s, "논문 1편 처리 비용 — 실측 (arXiv 2412.11427)", "본문은 12,000자로 절단 → 논문 길이와 무관하게 편당 비용 거의 고정")
pic(s, FIG+"c1_review.png", 0.4, 1.5, w=8.0)
box, tf = tb(s, 8.6, 1.7, 4.5, 5.0)
para(tf, "리뷰 실측", 14, RED, True, first=True, space_after=6)
runs(tf, [("입력 ", DARK, False), ("4,887", RED, True), (" tok", DARK, False)], 13, space_after=3)
runs(tf, [("출력 ", DARK, False), ("2,563", RED, True), (" tok (실제 생성)", DARK, False)], 13, space_after=3)
runs(tf, [("→ 리뷰 ", DARK, False), ("$0.0177", RED, True), (" (₩24)", GRAY, False)], 14, space_after=10)
para(tf, "9p · 49,225자 · 그림 2개", 11.5, GRAY, False, space_after=10)
para(tf, "1편 전체 합계", 14, RED, True, space_after=6)
runs(tf, [("리뷰+연결+그림+임베딩 ≈ ", DARK, False), ("$0.063", RED, True), (" (₩87)", GRAY, False)], 13, space_after=8)
para(tf, "동일 논문 재리뷰는 캐시 → 무료", 12, GREEN, True)
source(s, "측정: count_tokens(입력) + 실제 Haiku 생성 1회(출력) · WRITE_REVIEW_MODEL=claude-haiku-4-5, max_tokens=4000")

# ============================ SLIDE 4 — Deep Research ============================
s = slide()
title(s, "Deep Research 질의 비용 — 길이 × 모델", "검색된 ~8편 발췌(~7.5k tok) 기반 · 독자 BYOK 부담")
pic(s, FIG+"c2_deep.png", 0.4, 1.5, w=8.0)
box, tf = tb(s, 8.6, 1.7, 4.5, 5.0)
para(tf, "옵션", 14, BLUE, True, first=True, space_after=6)
para(tf, "Fast=Haiku · Smart=Sonnet · 길이 Short/Medium/Long", 12, DARK, False, bullet=("•",BLUE))
runs(tf, [("Deeper(Opus) ≈ ", DARK, False), ("$2.2", RED, True), (" — 약 70배", GRAY, False)], 12.5, space_after=8)
para(tf, "비용 부담", 14, BLUE, True, space_after=6)
para(tf, "답변 생성 = 독자 BYOK · 쿼리 임베딩=운영자(거의 0)", 12, DARK, False, bullet=("•",BLUE))
para(tf, "주의: 로컬(text.md 포함) 질의는 입력 ~85k tok로 급증", 12, AMBER, True, bullet=("⚠",AMBER))
source(s, "검색 인덱스 RAG(BM25+dense+RRF) · 출력은 길이 옵션 상한 · 배포 웹 기준(text.md 없음)")

# ============================ SLIDE 5 — ai4s 전체 ============================
s = slide()
title(s, "ai4s 전체 파이프라인 비용 — 실측 규모 (8 카테고리·1,178편)", "콜드 풀빌드 ≈ $51 (₩70,000)")
pic(s, FIG+"c3_ai4s.png", 0.4, 1.5, w=8.0)
box, tf = tb(s, 8.6, 1.7, 4.5, 5.0)
para(tf, "핵심 발견", 14, RED, True, first=True, space_after=6)
runs(tf, [("Paper connections = ", DARK, False), ("캐시 안 됨", RED, True)], 13, space_after=3)
runs(tf, [("입력 ", DARK, False), ("1.93M tok", RED, True), (" (실측)", GRAY, False)], 12.5, space_after=3)
para(tf, "매 배치가 타 카테고리 ~1,000편 목록 반복 포함", 11.5, GRAY, False, space_after=8)
runs(tf, [("→ curate 1회마다 ", DARK, False), ("connections ~$10", RED, True)], 13, space_after=3)
para(tf, "타임라인($6.8)보다 큼", 12, AMBER, True, space_after=8)
para(tf, "cross-cat insights는 캐시 → 입력 불변 시 무료", 12, GREEN, True)
source(s, "connections·insights 입력은 ai4s 실데이터로 count_tokens 측정 · 출력만 추정")

# ============================ SLIDE 6 — 논문 1편 추가 ============================
s = slide()
title(s, "논문 1편 ai4s 추가 — 경로에 따라 100배 차이", "비용은 ‘토픽 전체 connections를 다시 도느냐’에서 갈림")
chip(s, 0.7, 1.65, 5.6, 1.0, GREEN, "Paper Curio 플러그인", WHITE, 16, True)
box, tf = tb(s, 0.8, 2.85, 5.4, 3.4)
para(tf, "이 논문만: 리뷰+연결+그림+임베딩", 13, DARK, False, first=True, bullet=("•",GREEN))
para(tf, "토픽 전체 연결·인사이트·타임라인 미실행", 13, DARK, False, bullet=("•",GREEN))
runs(tf, [("≈ ", DARK, False), ("$0.06  (₩90)", GREEN, True)], 22, space_after=4)
chip(s, 7.05, 1.65, 5.6, 1.0, RED, "파이프라인 curate", WHITE, 16, True)
box2, tf2 = tb(s, 7.15, 2.85, 5.4, 3.4)
para(tf2, "신규 리뷰 + 그림", 13, DARK, False, first=True, bullet=("•",RED))
runs(tf2, [("connections ", DARK, False), ("전체 재실행", RED, True), (" (캐시X)", GRAY, False)], 13, space_after=4)
runs(tf2, [("≈ ", DARK, False), ("$10  (₩14,000)", RED, True)], 22, space_after=4)
para(tf2, "+ insights $0.7 / + timeline $1.15 (옵션)", 12, GRAY, False)
box3, tf3 = tb(s, 0.7, 6.2, 12, 0.7)
runs(tf3, [("💡 잦은 단건 추가는 플러그인 경로, 토픽 전체 connections 재생성은 모아서(주 1회).", NAVY, True)], 14, first=True)
source(s, "connections 미캐시로 curate 비용이 신규 편수와 무관하게 ~$10 고정 · 플러그인은 per-paper 연결만")

# ============================ SLIDE 7 — wiki concept ============================
s = slide()
title(s, "Karpathy LLM wiki 운영 — 두 가지 방법", "PDF를 어떻게 처리·저장·검색하는가의 차이 (장단점·한계는 다음 장)")
cpath = FIG+"concept.png"
if os.path.exists(cpath):
    # center the 16:9 concept image
    pic(s, cpath, 1.85, 1.55, w=9.6)
else:
    chip(s, 0.8, 1.7, 5.7, 3.6, RGBColor.from_string("EAF0FB"), "", BLUE)
    b, tf = tb(s, 1.0, 1.85, 5.3, 3.3)
    para(tf, "A. paper-curation", 16, BLUE, True, first=True, space_after=6)
    para(tf, "PDF → LLM 리뷰(distill) → 큐레이션 라이브러리(카테고리·연결·타임라인) → 짧은 발췌로 싼 질의", 13, DARK, False)
    para(tf, "선행 ↑ · 질의 ↓", 14, BLUE, True)
    chip(s, 6.85, 1.7, 5.7, 3.6, RGBColor.from_string("FDEEE3"), "", AMBER)
    b2, tf2 = tb(s, 7.05, 1.85, 5.3, 3.3)
    para(tf2, "B. PDF 쌓기 + RAG", 16, AMBER, True, first=True, space_after=6)
    para(tf2, "원문 PDF를 그대로 벡터 인덱싱 → 질의마다 큰 원문 chunk 회수", 13, DARK, False)
    para(tf2, "선행 ≈ 0 · 질의 ↑", 14, AMBER, True)
source(s, "개념도: PaperBanana 생성 · 모델: 논문 13k tok, 리뷰 $0.0177(실측), 답변 Sonnet 동일 가정")

# ============================ SLIDE 8 — 두 방식 장단점/한계 ============================
s = slide()
title(s, "두 방식의 장단점과 한계 — Karpathy LLM wiki", "같은 ‘위키’라도 상반된 강점·약점 (Karpathy: LLM = lossy·jagged)")
chip(s, 0.6, 1.5, 5.95, 0.6, BLUE, "A. paper-curation  (Curate → Query)", WHITE, 13.5, True)
bA, tfA = tb(s, 0.72, 2.28, 5.75, 4.1)
para(tfA, "장점", 12.5, GREEN, True, first=True, space_after=3)
for t in ["사람이 읽는 큐레이션 산출물 (리뷰·카테고리·연결·타임라인)",
          "독해 1회 amortize — 질의는 ‘사전 이해’ 재사용",
          "distill 고신호 + figure 사전추출 선택 표시 (답변 토큰↓)",
          "전역 구조·cross-paper 종합 (Deeper 그래프)·질의비 고정"]:
    para(tfA, t, 11, DARK, False, bullet=("✓", GREEN), space_after=2)
para(tfA, "한계 (Karpathy LLM)", 12.5, RED, True, space_after=3)
for t in ["압축 손실 — 리뷰에 없는 세부(수치·수식)는 접근 불가",
          "리뷰 모델 오류가 모든 질의에 전파 (품질 상한)",
          "선행 비용·지연 — 큐레이션 전엔 질의 불가",
          "고정 스키마 밖 질문 빈약 · 재큐레이션 비용"]:
    para(tfA, t, 11, DARK, False, bullet=("✗", RED), space_after=2)
chip(s, 6.78, 1.5, 5.95, 0.6, AMBER, "B. PDF 쌓기  (Pile → Query)", WHITE, 13.5, True)
bB, tfB = tb(s, 6.9, 2.28, 5.75, 4.1)
para(tfB, "장점", 12.5, GREEN, True, first=True, space_after=3)
for t in ["선행비 ≈ 0 · 즉시 질의 가능",
          "원문 직접 접근 → 디테일 보존 (손실 없음)",
          "질의 1~10회 영역 최저가",
          "셋업 단순"]:
    para(tfB, t, 11, DARK, False, bullet=("✓", GREEN), space_after=2)
para(tfB, "한계 (Karpathy LLM)", 12.5, RED, True, space_after=3)
for t in ["매 질의마다 원문 재투입 (독해 amortize 없음)·노이즈",
          "Figure: 페이지 이미지 멀티모달 과금 or figure 손실",
          "긴 논문일수록 토큰↑·청크 경계로 표/그림 분절 (A는 ~고정)",
          "전역구조 부재·검색 종속·컨텍스트 한계(~77편)→RAG 강제"]:
    para(tfB, t, 11, DARK, False, bullet=("✗", RED), space_after=2)
chip(s, 0.6, 6.48, 12.13, 0.5, LGRAY,
     "공통 한계 (Karpathy LLM): 환각·lossy recall · 컨텍스트 윈도우 천장 · 대규모 코퍼스는 RAG 불가피",
     DARK, 12, True)
source(s, "한계 정리: Karpathy ‘LLM=lossy compression·jagged intelligence’ 관점 + RAG vs long-context 트레이드오프")

# ============================ SLIDE 9 — wiki bars ============================
s = slide()
title(s, "wiki 운영비 — 논문수 × 질의(10회) 기준", "질의 1~10회 영역에서는 PDF 쌓기(B)가 전 구간 저렴")
pic(s, FIG+"c4_wikibars.png", 2.15, 1.55, w=9.0)
box, tf = tb(s, 1.4, 6.25, 11, 0.7)
runs(tf, [("1,000편·10회: A ", DARK, False), ("$18.8", RED, True), ("  vs  B ", DARK, False),
          ("$2.9", BLUE, True), ("  →  ", DARK, False), ("6.5배", RED, True),
          (" (리뷰 선행비가 회수되기 전)", GRAY, False)], 14, first=True, align=PP_ALIGN.CENTER)
source(s, "A=리뷰+임베딩 후 RAG / B=원문 임베딩 후 RAG · 답변 Sonnet 동일 · 출력 상쇄")

# ============================ SLIDE 9 — crossover ============================
s = slide()
title(s, "Trade-off 손익분기 — Q* ≈ 6.8 N / k", "k = 질의당 인용 문서수 · k↑ → 큐레이션 더 빨리 유리 (k=8이 기존 0.83N)")
pic(s, FIG+"c6_ksens.png", 0.4, 1.55, w=8.0)
table(s, [["k (인용/질의)", "Q*/N", "Q* @ N=1,000"],
          ["2", "3.4", "~3,400"], ["4", "1.7", "~1,700"],
          ["8", "0.85", "~850"], ["16", "0.43", "~430"], ["32", "0.21", "~210"]],
      8.65, 1.6, 4.35, 2.55, col_w=[1.55, 1.0, 1.8], header_fill=RED)
box, tf = tb(s, 8.65, 4.35, 4.4, 2.6)
para(tf, "k↑ → B/A 질의비 1.09→1.54× (B가 상대적으로 비싸짐)", 11, DARK, False, first=True, bullet=("•",RED), space_after=4)
para(tf, "Figure·긴 논문이면 실제 격차 더 큼 → Q* 더 ↓", 11, AMBER, True, bullet=("•",AMBER), space_after=4)
para(tf, "‘전부 넣기’ 불가: 1M÷13k ≈ 77편 → RAG 강제", 11, DARK, False, bullet=("•",GRAY), space_after=4)
para(tf, "모델 민감도: 강한 답변모델일수록 Q* ↓", 11, DARK, False, bullet=("•",NAVY))
source(s, "Q*=6.8N/k (k=8이 기존 0.83N) · A는 사전독해 amortize+figure 사전추출, B는 매 질의 원문 재투입")

# ============================ SLIDE 11 — 규모 상한 ============================
s = slide()
title(s, "Karpathy LLM wiki — 문헌 몇 편까지? (무제한 아님)",
      "방식별 천장이 다름 · 인덱스 실측 13 KB/편 (ai4s 1,178편 = 15.7MB)")
pic(s, FIG+"c7_scale.png", 0.4, 1.6, w=8.0)
table(s, [["편수", "다운로드", "상태"],
          ["1천", "13 MB", "여유"],
          ["1만", "134 MB", "가능"],
          ["3만", "401 MB", "부담"],
          ["10만", "1.3 GB", "한계"]],
      8.7, 1.7, 4.3, 2.3, col_w=[1.3, 1.6, 1.4], header_fill=NAVY)
box, tf = tb(s, 8.7, 4.25, 4.4, 2.6)
para(tf, "‘전부 context 덤프’ = 1M÷13k ≈ 77편이 천장", 11.5, RED, True, first=True, bullet=("•",RED), space_after=4)
para(tf, "현 브라우저 client 구조 실용 한계 ~1~3만 편", 11.5, AMBER, True, bullet=("•",AMBER), space_after=4)
para(tf, "server 벡터DB로 옮기면 사실상 무제한 (품질·예산이 한계)", 11.5, GREEN, True, bullet=("•",GREEN), space_after=4)
para(tf, "컨텍스트=작업기억(제한) · retrieval=외부기억(확장)", 11, DARK, False, bullet=("•",NAVY))
source(s, "인덱스 실측: _search_index.json + .bin = 13 KB/편 (ai4s 1,178편 15.7MB) · 클라이언트 다운로드 기준")

# ============================ SLIDE 12 — 에이전트 on-demand ============================
s = slide()
title(s, "에이전트 on-demand 읽기 — Obsidian + Claude Code + PDF 폴더",
      "색인 없이 그때그때 읽는 ‘4번째 방식’ (= paper-curation의 수동 버전)")
table(s, [["방식", "검색 방식", "질의당 읽기", "실용 한계"],
          ["전부 context 덤프", "없음 (다 넣음)", "전체", "~77편"],
          ["에이전트 on-demand (이 방식)", "파일명·grep·에이전트 추론", "관련 PDF 통째 몇 편", "~수십~수백"],
          ["임베딩 RAG · 브라우저", "벡터 유사도", "top-k 청크", "~1~3만"],
          ["임베딩 RAG · server DB", "벡터 유사도", "top-k 청크", "~수백만+"]],
      0.5, 1.55, 12.3, 2.25, col_w=[3.3, 3.0, 3.0, 3.0], header_fill=NAVY)
box, tf = tb(s, 0.6, 4.2, 7.1, 2.6)
para(tf, "메커니즘 — 임베딩·색인 없음", 12.5, RED, True, first=True, space_after=3)
para(tf, "Read·grep·glob로 PDF를 그때그때 펼침 (‘PDF 따로 안 건드림’ 맞음)", 11, DARK, False, bullet=("•",GRAY), space_after=3)
para(tf, "grep은 PDF 바이너리 본문 못 읽음 → 내용 검색이 약함", 11, DARK, False, bullet=("•",GRAY), space_after=3)
para(tf, "큰 폴더에선 ‘어느 PDF가 관련?’ 못 찾음 (파일명·폴더구조 의존)", 11, DARK, False, bullet=("•",GRAY), space_after=3)
para(tf, "질의당 원문 통째 → 많은 논문 가로지르는 질문은 토큰·context 부족", 11, DARK, False, bullet=("•",GRAY))
chip(s, 7.9, 4.2, 5.0, 2.55, RGBColor.from_string("EAF0FB"), "", BLUE)
b2, tf2 = tb(s, 8.1, 4.4, 4.6, 2.15)
para(tf2, "⭐ 핵심 — 사실상 paper-curation을 손으로", 12.5, NAVY, True, first=True, space_after=5)
para(tf2, "쌓이는 wiki .md = 점진적 수동 큐레이션 (grep으로 찾는 distill 층)", 11.5, DARK, False, space_after=5)
para(tf2, "시간이 지날수록 A(paper-curation)에 수렴 — A는 이를 배치+임베딩으로 자동·대규모화한 버전", 11.5, DARK, False)
source(s, "Claude Code = 에이전트 도구(Read/grep/glob)로 just-in-time 읽기 · 소규모(수십~수백)엔 최고로 간편, 대규모는 임베딩 인덱스 필요")

# ============================ SLIDE 13 — 결론 ============================
s = slide()
title(s, "결론 — 비용은 ‘싼 질의’가 아니라 ‘산출물’을 산다")
box, tf = tb(s, 0.7, 1.5, 12, 3.0)
runs(tf, [("순수 질의 비용", DARK, True), ("만 보면 1~10회 수준에선 ", DARK, False),
          ("PDF 쌓기+RAG가 압도적으로 쌈", BLUE, True), (" (리뷰 선행비 회수 불가).", DARK, False)], 15, first=True, space_after=10)
para(tf, "paper-curation의 비용이 사는 것:", 14, RED, True, space_after=5)
para(tf, "한글 리뷰·카테고리·연결·타임라인·네트워크 = 사람이 읽고 탐색하는 큐레이션 산출물", 13, DARK, False, bullet=("①",RED))
para(tf, "distill된 고신호 컨텍스트 (질의당 토큰 ↓)", 13, DARK, False, bullet=("②",RED))
para(tf, "Deeper 연결그래프 기반 멀티에이전트 리포트", 13, DARK, False, bullet=("③",RED))
# bottom compare chips
chip(s, 0.9, 4.95, 5.5, 1.2, RGBColor.from_string("EAF0FB"),
     "질의만 필요 → PDF + RAG (검색엔진)", BLUE, 15, True)
chip(s, 6.95, 4.95, 5.5, 1.2, RGBColor.from_string("FBEAEA"),
     "읽고 탐색·재사용 → paper-curation (저장소)", RED, 15, True)
box2, tf2 = tb(s, 0.9, 6.35, 11.6, 0.6)
runs(tf2, [("비용 절감 키: ", DARK, True),
           ("① 잦은 추가는 플러그인  ② connections 재생성은 모아서  ③ insights는 캐시 활용", GRAY, False)],
     12.5, first=True, align=PP_ALIGN.CENTER)
source(s, "전체 수치·표: operations/price/price.html · 측정 2026-06-19 · paper-curation 파이프라인 실측")

prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
